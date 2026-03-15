from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_prism_ppt():
    prs = Presentation()

    # Define some colors for the Prism theme
    PURPLE = RGBColor(129, 140, 248)
    CYAN = RGBColor(6, 182, 212)
    PINK = RGBColor(244, 114, 182)
    WHITE = RGBColor(255, 255, 255)
    DARK = RGBColor(13, 13, 20)

    def add_slide_with_background(layout_index):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        # Add a dark background shape
        background = slide.shapes.add_shape(
            1, # Rectangle
            0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = DARK
        background.line.width = 0
        return slide

    # Slide 1: Title
    slide = add_slide_with_background(6) # Blank layout
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    title_p = title_box.text_frame.add_paragraph()
    title_p.text = "PRISM"
    title_p.font.bold = True
    title_p.font.size = Pt(88)
    title_p.font.color.rgb = CYAN
    title_p.alignment = PP_ALIGN.CENTER

    sub_p = title_box.text_frame.add_paragraph()
    sub_p.text = "Early Cognitive Screening Through Play"
    sub_p.font.size = Pt(28)
    sub_p.font.color.rgb = WHITE
    sub_p.alignment = PP_ALIGN.CENTER

    # Slide 2: Team
    slide = add_slide_with_background(6)
    t_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    t_title_p = t_title.text_frame.add_paragraph()
    t_title_p.text = "Our Team"
    t_title_p.font.size = Pt(36)
    t_title_p.font.color.rgb = PURPLE

    names = ["AKANKSHA (IGDTUW CSE 2028)", "UDAY (DTU CSE 2028)", "VAIBHAV (DTU CSE 2028)", "VINEET (DTU CSE 2028)", "YASH (DTU CSE 2028)"]
    for i, name in enumerate(names):
        nb = slide.shapes.add_textbox(Inches(1), Inches(1.5 + i*1), Inches(8), Inches(1))
        np = nb.text_frame.add_paragraph()
        np.text = "• " + name
        np.font.size = Pt(24)
        np.font.color.rgb = WHITE

    # Slide 3: Problem
    slide = add_slide_with_background(6)
    prob_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    prob_title.text_frame.text = "PROBLEM STATEMENT"
    prob_title.text_frame.paragraphs[0].font.size = Pt(36)
    prob_title.text_frame.paragraphs[0].font.color.rgb = PINK

    p_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = p_content.text_frame
    tf.word_wrap = True
    p1 = tf.add_paragraph()
    p1.text = "Millions of neurodivergent children (Autism, ADHD, Dyslexia) face significant barriers in accessing personalized and engaging learning support. Parents struggle to monitor real-time development."
    p1.font.size = Pt(22)
    p1.font.color.rgb = WHITE

    stats = [
        "ADHD: 5-15% symptomatic in children",
        "Autism (ASD): 1-3% of children",
        "Dyslexia: 6-12% of children",
        "Cerebral Palsy: 2.1-3 per 1,000 live births"
    ]
    for s in stats:
        sp = tf.add_paragraph()
        sp.text = "• " + s
        sp.font.size = Pt(18)
        sp.font.color.rgb = CYAN

    # Slide 4: Solution Overview
    slide = add_slide_with_background(6)
    sol_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    sol_title.text_frame.text = "THE PRISM SOLUTION"
    sol_title.text_frame.paragraphs[0].font.size = Pt(36)
    sol_title.text_frame.paragraphs[0].font.color.rgb = CYAN

    p_a_c = ["PLAY (Gamified Experience)", "ANALYZE (AI Reporting)", "CONNECT (Community Care)"]
    for i, txt in enumerate(p_a_c):
        tb = slide.shapes.add_textbox(Inches(0.5 + i*3.2), Inches(1.5), Inches(3), Inches(1))
        tp = tb.text_frame.add_paragraph()
        tp.text = txt
        tp.font.size = Pt(18)
        tp.font.bold = True
        tp.font.color.rgb = PURPLE

    desc = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(3))
    dt = desc.text_frame
    dp = dt.add_paragraph()
    dp.text = "We move beyond simple Yes/No diagnosis to actionable data using 7 specialized cognitive games and AI reporting via 'PRISMy' assistant."
    dp.font.size = Pt(22)
    dp.font.color.rgb = WHITE

    # Slide 5: ADHD Detection
    slide = add_slide_with_background(6)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1)).text_frame.text = "ADHD: Attention & Impulse"
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
    p = tf.add_paragraph()
    p.text = "🎯 Attention Check: Measuring focus, distraction, and reaction consistency."
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "🚦 Stop Signal: Gold standard for testing the brain's 'brakes' (Response Inhibition)."
    p.font.size = Pt(22)
    p.font.color.rgb = CYAN

    # Slide 6: Autism Detection
    slide = add_slide_with_background(6)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1)).text_frame.text = "AUTISM: Social & Motor"
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
    p = tf.add_paragraph()
    p.text = "😊 Emotion Lab: Assessing facial emotional recognition and empathy cues."
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "🏃 Motor Match: CAMI-based motor imitation testing (80%+ detection accuracy)."
    p.font.size = Pt(22)
    p.font.color.rgb = PINK

    # Slide 7: Dyslexia Detection
    slide = add_slide_with_background(6)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1)).text_frame.text = "DYSLEXIA: Spatial & Auditory"
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
    p = tf.add_paragraph()
    p.text = "🧭 Direction Sense: Testing mirror confusion patterns and spatial lag."
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "🎵 Sound Pattern: Phonological processing through musical sequence recall."
    p.font.size = Pt(22)
    p.font.color.rgb = PURPLE

    # Slide 8: Memory & AI assistant
    slide = add_slide_with_background(6)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1)).text_frame.text = "Memory & AI Support"
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
    p = tf.add_paragraph()
    p.text = "🧩 Memory Matrix: Working memory span testing — backed by NIH research."
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "🤖 PRISMy Assistant: Real-time AI chat for psychoeducation and tips."
    p.font.size = Pt(22)
    p.font.color.rgb = CYAN

    # Slide 9: Results Dashboard
    slide = add_slide_with_background(6)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1)).text_frame.text = "PROBABILISTIC RISK REPORT"
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
    p = tf.add_paragraph()
    p.text = "• AI-Driven Analysis of error rates and response times"
    p.text += "\n• 7-Domain Cognitive Profile Chart"
    p.text += "\n• Actionable AI Recommendations for parents/therapists"
    p.text += "\n• Longitudinal tracking of developmental growth"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE

    # Slide 10: Business Model
    slide = add_slide_with_background(6)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1)).text_frame.text = "BUSINESS MODEL"
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
    p = tf.add_paragraph()
    p.text = "Freemium Tier: First 3 reports are free to build trust."
    p.font.size = Pt(20)
    p.font.color.rgb = CYAN
    p = tf.add_paragraph()
    p.text = "Premium: Longitudinal progress tracking & analytics."
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "Partner Model: Therapists pay platform fee for digital tools."
    p.font.size = Pt(20)
    p.font.color.rgb = PURPLE

    # Slide 11: Impact & Tech
    slide = add_slide_with_background(6)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1)).text_frame.text = "IMPACT & TECH STACK"
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
    p = tf.add_paragraph()
    p.text = "Impact: Early screening prevents academic decline & improves regulation."
    p.font.size = Pt(20)
    p.font.color.rgb = PINK
    p = tf.add_paragraph()
    p.text = "Tech Stack: HTML5, CSS3, JavaScript (Core Platform), Web Audio API."
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "Live Demo: https://bingo5706.vercel.app/"
    p.font.size = Pt(20)
    p.font.color.rgb = CYAN

    prs.save('PRISM_Presentation.pptx')
    print("PPT generated: PRISM_Presentation.pptx")

if __name__ == "__main__":
    create_prism_ppt()
